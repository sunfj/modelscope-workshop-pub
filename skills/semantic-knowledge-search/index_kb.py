#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
知识库索引脚本 - 供 install.py 调用，也可独立运行。
扫描 knowledge/ 目录，加载 BGE 模型生成向量并保存到 index_store/。
"""

import os
import sys
import gc
import sqlite3
import numpy as np
from pathlib import Path

os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"
os.environ["HF_HUB_DISABLE_TELEMETRY"] = "1"

# ---- 默认路径 ----
SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_ROOT = SCRIPT_DIR.parent.parent
DEFAULT_KNOWLEDGE_DIR = PROJECT_ROOT / "lab5-local-knowledge-assistant" / "knowledge"
DEFAULT_INDEX_STORE_DIR = PROJECT_ROOT / "lab5-local-knowledge-assistant" / "index_store"
DEFAULT_BGE_MODEL_DIR = PROJECT_ROOT / "lab5-local-knowledge-assistant" / "models" / "bge-small-zh-v1.5"


def _resolve_paths(knowledge_dir: Path = None, bge_model_dir: Path = None):
    """Resolve paths. If knowledge_dir is given, index_store and bge_model are derived from it."""
    if knowledge_dir is None:
        knowledge_dir = DEFAULT_KNOWLEDGE_DIR
    if bge_model_dir is None:
        bge_model_dir = DEFAULT_BGE_MODEL_DIR
    # index_store 放在知识库所在目录的同级
    index_store_dir = knowledge_dir.parent / "index_store"
    return knowledge_dir, index_store_dir, bge_model_dir


def _find_vl_model():
    """Find Qwen3-VL model in common locations."""
    candidates = [
        SCRIPT_DIR.parent.parent / "lab1-multimodal-vlm" / "Qwen3-VL-4B-Instruct-int4-ov",
        SCRIPT_DIR.parent.parent / "lab1-multimodal-vlm" / "Qwen3-VL-4B-OpenVINO",
        SCRIPT_DIR.parent.parent / "Qwen3-VL-4B-Instruct-int4-ov",
        SCRIPT_DIR.parent.parent / "lab5-local-knowledge-assistant" / "Qwen3-VL-4B-Instruct-int4-ov",
        Path("Qwen3-VL-4B-Instruct-int4-ov"),
        Path("lab1-multimodal-vlm") / "Qwen3-VL-4B-Instruct-int4-ov",
    ]
    for c in candidates:
        if c.exists():
            # Check for OpenVINO model files
            for model_file in ["openvino_model.xml", "openvino_model.bin", "config.json"]:
                if (c / model_file).exists():
                    return c
    return None


def download_model_if_missing(bge_model_dir: Path):
    if bge_model_dir.exists() and (bge_model_dir / "config.json").exists():
        return
    print(f"📥 语义模型未找到，从 ModelScope 下载...")
    try:
        from modelscope import snapshot_download
        snapshot_download(
            "AI-ModelScope/bge-small-zh-v1.5",
            local_dir=str(bge_model_dir),
        )
        print("✅ BGE 模型下载完成")
    except Exception as e:
        print(f"❌ 模型下载失败: {e}")
        sys.exit(1)


def parse_file(path: Path) -> str:
    """Extract text from a file based on its extension."""
    suffix = path.suffix.lower()
    try:
        if suffix in (".txt", ".md", ".log", ".nfo"):
            return path.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".pdf":
            import pypdf
            return "\n".join(
                page.extract_text() or "" for page in pypdf.PdfReader(path).pages
            )
        elif suffix == ".docx":
            import docx
            return "\n".join(p.text for p in docx.Document(path).paragraphs)
        elif suffix == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            return "\n".join(
                str(c)
                for ws in wb.worksheets
                for r in ws.iter_rows(values_only=True)
                for c in r if c
            )
        elif suffix == ".pptx":
            import pptx
            return "\n".join(
                p.text
                for s in pptx.Presentation(path).slides
                for sh in s.shapes
                if sh.has_text_frame
                for p in sh.text_frame.paragraphs
            )
        elif suffix in (".png", ".jpg", ".jpeg"):
            # OCR mode - load Qwen3-VL if available
            return ocr_image(path)
    except Exception as e:
        print(f"⚠️ 解析失败 {path.name}: {e}")
        return ""
    return ""


def ocr_image(img_path: Path) -> str:
    """Use Qwen3-VL to OCR an image. Skip if model not available."""
    model_dir = _find_vl_model()

    if model_dir is None:
        print(f"  ⏭️ 跳过 OCR (Qwen3-VL 模型不存在): {img_path.name}")
        print(f"  💡 支持的模型目录名: Qwen3-VL-4B-Instruct-int4-ov, Qwen3-VL-4B-OpenVINO")
        return ""

    print(f"  🤖 使用模型: {model_dir}")

    try:
        from optimum.intel.openvino import OVModelForVisualCausalLM
        from transformers import AutoProcessor
        from PIL import Image

        # Thumbnail for speed
        thumb = img_path.with_suffix(".thumb.jpg")
        with Image.open(img_path) as img:
            img = img.convert("RGB")
            img.thumbnail((256, 256), Image.Resampling.LANCZOS)
            img.save(thumb, quality=85)

        print(f"  👁️ 正在 OCR: {img_path.name}...")
        model = OVModelForVisualCausalLM.from_pretrained(str(model_dir), device="CPU")
        processor = AutoProcessor.from_pretrained(str(model_dir))

        msgs = [{
            "role": "user",
            "content": [
                {"type": "image", "image": str(thumb)},
                {"type": "text", "text": "请提取图片中的所有文字和主要内容。"},
            ],
        }]
        inputs = processor.apply_chat_template(
            msgs, tokenize=True, add_generation_prompt=True,
            return_dict=True, return_tensors="pt",
        )
        out = model.generate(**inputs, max_new_tokens=256)
        text = processor.decode(out[0], skip_special_tokens=True).split("assistant")[-1].strip()

        del model, processor
        gc.collect()
        if thumb.exists():
            thumb.unlink()

        print(f"  ✅ OCR 完成: {img_path.name} -> {text[:50]}...")
        return text
    except Exception as e:
        print(f"  ⚠️ OCR 失败 {img_path.name}: {e}")
        return ""


def build_index(knowledge_dir: Path = None, bge_model_dir: Path = None, include_ocr: bool = False):
    """Scan knowledge dir, extract text, embed, save."""
    knowledge_dir, index_store_dir, bge_model_dir = _resolve_paths(knowledge_dir, bge_model_dir)

    if not knowledge_dir.exists():
        print(f"❌ 知识库目录不存在: {knowledge_dir}")
        print("💡 请将文件放入后重试")
        sys.exit(1)

    index_store_dir.mkdir(parents=True, exist_ok=True)
    db_path = index_store_dir / "knowledge.db"
    emb_path = index_store_dir / "embeddings.npy"

    # Init DB
    with sqlite3.connect(db_path) as conn:
        conn.execute(
            "CREATE TABLE IF NOT EXISTS chunks "
            "(id INTEGER PRIMARY KEY, file_path TEXT, file_type TEXT, chunk_text TEXT, file_name TEXT)"
        )

    # Download model if missing
    download_model_if_missing(bge_model_dir)

    # Load BGE
    print(f"🧠 加载 BGE 语义模型...")
    from sentence_transformers import SentenceTransformer
    model = SentenceTransformer(str(bge_model_dir), device="cpu")

    # Scan files
    files = sorted(knowledge_dir.iterdir())
    for f in files:
        if not f.is_file():
            continue
        if f.suffix.lower() == ".thumb.jpg":
            continue
        if f.suffix.lower() in (".png", ".jpg", ".jpeg") and not include_ocr:
            print(f"  ⏭️ 跳过图片 (无 OCR): {f.name}")
            continue

        print(f"  📄 解析: {f.name}")
        text = parse_file(f)
        if not text.strip():
            continue

        # Chunk
        pieces = [text[i:i + 500] for i in range(0, len(text), 500)]
        with sqlite3.connect(db_path) as conn:
            conn.execute("DELETE FROM chunks WHERE file_path = ?", (str(f),))
            conn.executemany(
                "INSERT INTO chunks (file_path, file_type, chunk_text, file_name) VALUES (?, ?, ?, ?)",
                [(str(f), f.suffix.lower(), c, f.name) for c in pieces],
            )
            conn.commit()
        print(f"  ✅ 入库: {f.name} ({len(pieces)} 块)")

    # Load all chunks and embed
    with sqlite3.connect(db_path) as conn:
        all_chunks = conn.execute(
            "SELECT id, file_path, file_type, chunk_text, file_name FROM chunks"
        ).fetchall()

    if not all_chunks:
        print("⚠️ 知识库为空，无文件可索引")
        return

    print(f"🔄 生成语义向量 ({len(all_chunks)} 块)...")
    texts = [c[3] for c in all_chunks]
    embeddings = model.encode(texts, normalize_embeddings=True, batch_size=8)
    np.save(emb_path, embeddings)

    print(f"✅ 索引完成: {len(all_chunks)} 个文本块 -> {emb_path}")

    # Cleanup
    del model
    gc.collect()


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="索引知识库")
    parser.add_argument("--knowledge-dir", type=Path, default=None, help="知识库目录路径")
    parser.add_argument("--bge-model-dir", type=Path, default=None, help="BGE 模型路径")
    parser.add_argument("--with-ocr", action="store_true", help="包含图片 OCR (需要 Qwen3-VL 模型)")
    args = parser.parse_args()
    build_index(
        knowledge_dir=args.knowledge_dir,
        bge_model_dir=args.bge_model_dir,
        include_ocr=args.with_ocr,
    )
